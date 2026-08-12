"""Build the project-journey slide deck (pptx).

Covers: initial hypothesis, setting, original findings, feedback, 2026
landscape, controls, decisive within-problem experiment, length confound,
probe baseline, conclusions, open work.
"""

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

INK = RGBColor(0x0B, 0x0B, 0x0B)
SUB = RGBColor(0x52, 0x51, 0x4E)
BLUE = RGBColor(0x2A, 0x78, 0xD6)
ORANGE = RGBColor(0xEB, 0x68, 0x34)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_slide(title, bullets=None, image=None, img_width=None, notes=None,
              title_color=INK):
    s = prs.slides.add_slide(BLANK)
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = title_color
    top = 1.35
    if bullets:
        bb = s.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12.1),
                                  Inches(5.8 if not image else 2.4))
        tf = bb.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            if b.startswith("  "):
                para.level = 1
                b = b.strip()
            para.text = b
            para.font.size = Pt(16 if len(bullets) > 6 else 18)
            para.font.color.rgb = INK
            para.space_after = Pt(6)
        top += 2.5 if image else 0
    if image and os.path.exists(image):
        w = Inches(img_width or 8.5)
        s.shapes.add_picture(image, Inches(0.7), Inches(top), width=w)
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    return s


# 1 title
s = prs.slides.add_slide(BLANK)
tb = s.shapes.add_textbox(Inches(1), Inches(2.4), Inches(11.3), Inches(2.5))
tf = tb.text_frame
tf.paragraphs[0].text = "Attention Patterns and CoT Faithfulness"
tf.paragraphs[0].font.size = Pt(44)
tf.paragraphs[0].font.bold = True
p = tf.add_paragraph()
p.text = "A project journey: hypothesis -> feedback -> controls -> answer"
p.font.size = Pt(24)
p.font.color.rgb = SUB
p = tf.add_paragraph()
p.text = "Yulia Volkova | Athena Fellowship -> MATS 9.1 | Nov 2025 - Aug 2026"
p.font.size = Pt(18)
p.font.color.rgb = SUB

# 2 hypothesis
add_slide("Initial hypothesis (Nov 2025)", [
    "If attention attribution reflects how models use information, then:",
    "  Faithful CoT (verbalizes the hint it follows) -> focused attention on key sentences",
    "  Unfaithful CoT (follows silently, post-hoc rationalization) -> diffuse attention",
    "Operationalized via Thought Anchors (Bogdan et al. 2025): receiver heads =",
    "  attention heads with high kurtosis over sentence-level received attention",
    "If true: a white-box, unsupervised faithfulness monitor for deployed models",
])

# 3 setting
add_slide("Setting & pipeline",
          ["Model: DeepSeek-R1-Distill-Qwen-14B | Datasets: MMLU (143), GPQA-Diamond (186)",
           "Cue: 'A Stanford professor thinks the answer is (X)' with X != ground truth",
           "20 rollouts per problem -> label: verbalizing (mentions cue) vs silent",
           "Problem filters: low no-reasoning accuracy, cue-response gap >= 0.5",
           "Attention: sentence-level vert scores -> per-head kurtosis -> receiver heads"],
          image="methodology_flowchart.png", img_width=7.5)

# 4 original findings
add_slide("Original findings (Jan 2026, thought-anchors-faithfulness.vercel.app)", [
    "1. Faithful CoT showed higher kurtosis (GPQA p=0.04, rollout-level)",
    "2. Reasoning-only receiver heads barely overlap between groups (n.s.)",
    "3. Unfaithful CoT followed the cue more (80-91% vs 66-71%)",
    "4. Universal receiver heads across datasets and conditions (p<0.001)",
    "Interactive webapp for exploring attention per problem",
])

# 5 feedback
add_slide("Feedback round (Feb-Mar 2026)", [
    "Adam Karvonen: bigger n, second model family, p<0.01; kurtosis-as-unsupervised-",
    "  metric is the contribution, supervised probes are easy; optimize the pipeline",
    "Goodfire (Jack Merullo, Sid Boppana): killer app = unknown-cue detection;",
    "  three-legged evidence (kurtosis + probe + attention blocking); FlashAttention caveat",
    "Uzay Macar: probes with kurtosis as unsupervised metric; precise research question",
    "Daria & Riya: verify cue presence is not the only source of receiver heads",
])

# 6 landscape
add_slide("Meanwhile, the field moved (Dec 2025 - Jul 2026)", [
    "Zaman & Srivastava (ACL 2026): hint-verbalization is a flawed faithfulness label;",
    "  non-verbalized hints still causally mediate through the CoT; faithful@k",
    "Arcuschin et al.: black-box pipeline finds unverbalized biases (~$100/bias,",
    "  population-level only)",
    "Rationalization probes (2603.17199): supervised activation probes detect",
    "  motivated reasoning, AUC 65-82%, even before CoT generation",
    "-> reframe: the open question is the UNSUPERVISED attention signal",
])

# 7 control 1
add_slide("Control 1: cue-conditioning (Daria & Riya's check)", [
    "Cue sentence was never in the metric (drop_first ate position 0) -> cue salience",
    "  cannot explain the kurtosis gap; but position 0 = attention sink (design flaw)",
    "GPQA: faithful/unfaithful gap persists in UNCUED rollouts (d~1.2)",
    "  -> gap tracks problem identity, not faithful behavior",
    "PI-level stats: original p=0.04 was a rollout-level artifact (best honest p=0.056)",
])

# 8 control 2
add_slide("Control 2: is the anchor just the cue restatement?",
          ["Cue-mention sentences are usually NOT the top-attended sentences",
           "Removing them changes kurtosis by <1% (points on the diagonal)",
           "Dominant spike = first reasoning sentence, not the cue"],
          image="plots/cue_ablation/ablation_scatter.png", img_width=4.6)

# 9 decisive design
add_slide("The decisive experiment: within-problem comparison", [
    "Confound-proof design: compare verbalizing vs silent rollouts of the SAME problem",
    "76 mixed problems (42 MMLU, 34 GPQA), 3+3 rollouts each = 456 rollouts",
    "Teacher-forced re-extraction through the generating 14B on one H100",
    "New pipeline: attention pooled to sentence level inside the attention kernel,",
    "  in 8-head chunks -> 14B x 8k tokens feasible (~45 GB peak saved)",
    "Same forward passes also save activations -> supervised probe baseline for free",
])

# 10 result + confound
add_slide("Result: a 'signal' appears - and dissolves",
          ["Paired within-problem kurtosis: d=0.42, p<0.001, 74% of problems positive...",
           "...but kurtosis correlates rho=0.98 with rollout length,",
           "and verbalizing rollouts are +19 sentences longer (same 76% positive)"],
          image="plots/final/kurt_vs_length.png", img_width=6.2)

# 11 length controls
add_slide("Length controls: the signal is length",
          ["Length-residualized paired test: p=0.53, 46% positive (coin flip)",
           "Length-matched pairs (|dn|<=5): concordance 0.549 ~ chance (133 pairs)"],
          image="plots/final/paired_diffs.png", img_width=8.2)

# 12 probe ceiling
add_slide("The supervised ceiling: internals do contain the signal",
          ["Activation probe (GroupKFold by problem): post-CoT AUC 0.83, pre-CoT 0.64",
           "Kurtosis within-problem AUC 0.65 - and adds NOTHING on top of the probe",
           "Detection of unverbalized cue use is possible - but needs supervision"],
          image="plots/final/auc_bars.png", img_width=6.8)

# 13 conclusions
add_slide("Conclusions", [
    "The hypothesized unsupervised attention signature of (un)faithful CoT does not",
    "  exist once length and problem identity are controlled",
    "Earlier positive results (incl. mine) were length + problem-selection artifacts",
    "Positive findings:",
    "  Verbalization costs ~19 sentences (supports the incompleteness account)",
    "  Pre-CoT activations predict upcoming verbalization (AUC 0.64)",
    "  Memory-safe receiver-head extraction at 14B/8k scale (reusable pipeline)",
    "  Hiring setting: strong INVERTED targeted-attention signal (next slides)",
], title_color=BLUE)

# 14 entropy verdict
add_slide("The diffuseness test: entropy (strict criteria, 25 problems)", [
    "Normalized attention entropy (length-invariant by design), 1,556 rollouts extracted;",
    "  analysis on the 25 problems meeting ALL original criteria (383 rollouts)",
    "Raw: verbalizing MORE diffuse (opposite sign, p=0.03) - but entropy still tracks",
    "  length (rho 0.49-0.72), and verbalizing rollouts are longer",
    "Length-residualized: d ~ 0.0 (p=0.73) | length-matched pairs: P = 0.516/0.488 (chance)",
    "Verdict: NO diffuseness difference in either direction - every global attention",
    "  statistic inherits the verbalization-length coupling",
])

# 15 hiring result
add_slide("Hiring setting (Karvonen & Marks): the inverted positive result", [
    "Targeted test: attention from the decision token to the candidate's NAME",
    "  (gemma-3-12b-it, 111 resumes x 4 names x 4 prompts; ground truth = delta p(yes))",
    "The model attends to the name LESS when the name causally drives the decision:",
    "  resume-level Spearman -0.64 (p=4e-14); -0.57 controlling decision confidence;",
    "  length NOT a confound (rho -0.04); detector AUC ~0.87 (inverted)",
    "Interpretations: covert absorption during prefill vs learned attention suppression",
    "Usable risk signal: low overt attention to a sensitive attribute + borderline",
    "  decision -> elevated risk of covert influence",
], title_color=ORANGE, image="plots/final/hiring_scatter.png", img_width=5.2)

# 16 open
add_slide("Open threads", [
    "Distinguish absorption vs suppression: name-swap ablation mid-forward",
    "Replicate inverted attention result on Gemma-3-27B / Mistral-Small (data ready)",
    "Design fixes for any CoT successor: cue mid-prompt, generous token budget,",
    "  within-problem comparisons only",
    "Write-up: null (global statistics) + positive (targeted attention) + probe baselines",
])

prs.save("thought_anchors_journey.pptx")
print("saved thought_anchors_journey.pptx,", len(prs.slides.slides if hasattr(prs.slides,'slides') else prs.slides._sldIdLst), "slides")
